#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract the full text of a tender, deterministically, in about a minute.

    python3 normalize.py --in out/ --out out/normalized/

Three properties are load-bearing, in this order.

NOTHING USEFUL MAY VANISH. The rule that makes this provable is that usefulness is never
judged. Judging content would need a model; instead each file is classified by whether a
deterministic decoder can recover characters from it. Everything readable is extracted in
full and never filtered. Everything unreadable is listed by name, size and digest, so a gap
is visible rather than silent. No file is ever dropped on a guess about its importance.

SECONDS FOR A NORMAL TENDER, MINUTES FOR THE WORST CASE. A median tender is a handful of
DOCX and PDF and takes seconds. The worst case is a building project shipped as one large
archive — hundreds of files, most of them PDF — and it runs end to end in a couple of
minutes, most of that PDF text extraction and archive unpacking.

DETERMINISTIC, NO MODEL. Same bytes in, same text out. No LLM, and no OCR on this path:
tender PDFs carry a text layer, so OCR would buy nothing here while making the output
depend on an engine's version. Dependencies are pinned exactly (requirements.txt) and walk
orders are sorted, because "deterministic" must hold across installs and across time, not
merely within one lucky machine.

WHY PDF TABLE DETECTION IS GONE. It was the single heuristic in the extraction path, and it
earned its removal three times over: it cost orders of magnitude more time; it recovered no
words that plain extraction had missed, because table text is already in the text layer;
and on the rotated text in drawing title blocks it emitted characters backwards, which
plain extraction reads correctly. What detection added was grid layout, at great cost, with
corruption. Tables in DOCX/XLSX/PPTX are still extracted: those formats declare their
structure, so reading it is exact and free.

WHY PDF IMAGES ARE OFF BY DEFAULT. A text extract cannot use a raster, and extracting them
adds a great many megabytes and not one character. `--with-images` restores them.

ARCHIVES ARE UNPACKED, WITH LIMITS. Buyers ship the specification inside a ZIP or a 7z, and
those nest — an archive holding an archive holding the project. Treating one as opaque
leaves a tender looking documented and unread at once, the worst state because it is
invisible. Depth, file count and expanded size are capped so a bomb fails loudly instead of
filling the runner.
"""

import argparse
import hashlib
import io
import json
import os
import re
import shutil
import time
import zipfile

# Signed containers nest: a 7z holds a ZIP, the ZIP holds a signed container, and that one
# holds another. At depth 3 three real documents went unread. Expansion is still bounded by
# member count and byte budget, so raising this cannot let an archive bomb win.
MAX_DEPTH = 6
MAX_MEMBERS = 5000
MAX_EXPANDED = 8 * 1024 ** 3          # 8 GB — a runner has ~14 GB free
MAX_MEMBER_BYTES = 2 * 1024 ** 3

# Images are off the default path (see the module docstring). When --with-images is asked
# for, a picture counts only if a human would see it as one rather than as a rule, a logo
# fragment or a scanning artefact. Fixed on purpose: no heuristics that drift.
WITH_IMAGES = False
# Kept only when a later stage will read the files an archive gave up — see the gap
# list below. Off by default because the scratch can hold gigabytes.
KEEP_UNPACKED = False
MIN_IMAGE_PX = 64
MIN_IMAGE_PAGE_FRACTION = 0.01

# A file yields text or it does not, and that is the only judgement made about it. Below
# this many characters there is nothing a reader could act on, and the file is reported as
# unreadable with its name, size and digest rather than quietly counted as processed.
MIN_USEFUL_CHARS = 1

TEXTUAL = {".txt": "text", ".csv": "csv", ".tsv": "tsv", ".xml": "xml",
           ".json": "json", ".md": "markdown"}
ARCHIVES = {".zip", ".7z"}

# The OpenDocument family, recognised by the media type the format itself stores rather than
# by the extension a buyer typed. Estonian authorities publish price tables and specifications
# as .ods and .odt often enough that leaving them out is a real gap rather than a tidy edge
# case, and they are read exactly as Word 97 files are: converted once, without a model and
# without a network call.
ODF_MIMETYPE = "mimetype"
ODF = {
    "application/vnd.oasis.opendocument.spreadsheet": ".ods",
    "application/vnd.oasis.opendocument.text": ".odt",
    "application/vnd.oasis.opendocument.presentation": ".odp",
}


# ------------------------------------------------------------------ long paths on Windows
#
# WHY THIS EXISTS, MEASURED RATHER THAN ANTICIPATED. A four-day trial run against the live
# register lost 15 of 79 procurements to `WinError 206`: the path was too long. Not one of
# them failed to download, and not one of them was unreadable — the extraction simply could
# not create the file it was about to write, and the tender was reported as lost.
#
# The cause is the shape of the documents rather than a bug: a building project arrives as an
# archive of archives, and the extractor mirrors that nesting under `normalized/`. Add a
# runtime root, a tender id, a section and two levels of Estonian folder names and 260
# characters are gone. The register is not doing anything unusual; Windows is.
#
# The production runner is Linux, where the limit is 4096 and none of this fires. That is
# exactly why it is worth fixing rather than living with: the tool promises that a VPS, a
# laptop and a runner execute the identical thing, and a laptop that silently drops a fifth of
# a window does not. The extended-length prefix lifts the limit to 32,767 characters.
#
# TWO RULES COME WITH THE PREFIX. The path must be absolute, and it must use backslashes
# only: a prefixed path with a forward slash anywhere in it is not the file you meant, it
# is a file that does not exist. Member names inside an archive are joined with forward
# slashes on purpose -- they are the document's logical address and travel into the
# manifest -- so the separator has to be fixed here, at the moment a logical address
# becomes a real one.
_LONG = "\\\\?\\"


def fspath(path):
    """An absolute path the filesystem will accept however deep it is."""
    real = os.path.normpath(os.path.abspath(path))
    # A UNC path already begins with two backslashes and takes a different prefix. Left
    # alone rather than guessed at, because nothing in this tool writes to one.
    if os.name != "nt" or real.startswith(_LONG) or real.startswith("\\\\"):
        return real
    return _LONG + real


def plainpath(path):
    """The same path without the prefix, for handing to a program that may not know it.

    LibreOffice and 7z are given file names on a command line, and the prefix is
    understood by some tools and not by others. Inside this process it is what makes a
    deep tree writable; outside it, it is a risk with no benefit, because the binaries are
    only ever pointed at a scratch directory that is short by construction.
    """
    text = str(path)
    return text[len(_LONG):] if text.startswith(_LONG) else text


def slug(name, limit=80):
    """Filesystem-safe name that stays unique under truncation.

    Plain truncation destroyed a file: two archive members sharing their first 80
    characters converged on one output path, the second silently overwrote the first, and
    nothing appeared in the gap list — the one failure mode the audit trail cannot see.
    Truncation therefore appends a digest of the full name, and the extension survives
    because it is sniff()'s last resort for files without a magic signature.
    """
    s = re.sub(r"[^\w\-. ]+", "_", name, flags=re.UNICODE).strip(" ._") or "file"
    if len(s) <= limit:
        return s
    stem, ext = os.path.splitext(s)
    ext = ext[:12]
    digest = hashlib.sha1(s.encode("utf-8")).hexdigest()[:8]
    return stem[:max(1, limit - len(ext) - 9)] + "~" + digest + ext


def sniff(path):
    """What the file is, from its bytes. Buyers name attachments freely; content does not lie."""
    try:
        with open(path, "rb") as fh:
            head = fh.read(16)
    except OSError:
        return ""
    if head[:4] == b"%PDF":
        return ".pdf"
    if head[:12] == b"ISO-10303-21":
        return ".ifc"
    if head[:2] == b"7z" or head[:6] == b"7z\xbc\xaf\x27\x1c":
        return ".7z"
    if head[:4] == b"PK\x03\x04":
        try:
            with zipfile.ZipFile(path) as z:
                names = z.namelist()
                # Read while the archive is still open; the checks below run after it closes.
                media = (z.read(ODF_MIMETYPE).decode("ascii", "replace").strip()
                         if ODF_MIMETYPE in names else "")
            if "word/document.xml" in names:
                return ".docx"
            if any(n.startswith("xl/") for n in names):
                return ".xlsx"
            if any(n.startswith("ppt/") for n in names):
                return ".pptx"
            # OPEN-FORMAT OFFICE DOCUMENTS, AND WHY THEY NEED ASKING FOR BY NAME. An .ods
            # spreadsheet is a ZIP with none of the members above, so without this it falls
            # through as an archive and is UNPACKED: a reader is handed content.xml and
            # styles.xml instead of the buyer's price table, and nothing reports a failure.
            # The media type is the format's own answer, stored as the mandatory first
            # member. Signed containers carry a `mimetype` too — `…asic-e+zip` — and must
            # keep falling through to the archive path, which is why this matches the
            # OpenDocument family rather than the presence of the member.
            if media in ODF:
                return ODF[media]
        except (zipfile.BadZipFile, OSError, KeyError):
            pass
        return ".zip"
    if head[:5] == b"{\\rtf":
        return ".rtf"
    if head[:4] == b"\xd0\xcf\x11\xe0":
        return _ole_subtype(path)
    # `mimetype` is the mandatory first member of ASiC/ODF containers (.edoc among them):
    # a short plain-text media type. Twenty-five of them were flooding the gap list as
    # "unknown", which is noise exactly where silence must be meaningful.
    if os.path.basename(path) == "mimetype" and head[:16].isascii():
        return ".txt"
    return os.path.splitext(path)[1].lower().strip(")]}, ")


def _ole_subtype(path):
    """Word, Excel or PowerPoint? The OLE2 header is identical for all three, so the answer
    comes from the stream names in the directory, which are stored as UTF-16LE."""
    try:
        with open(path, "rb") as fh:
            head = fh.read(1 << 20)
    except OSError:
        return ".doc"
    def utf16(s):
        return s.encode("utf-16-le")
    if utf16("WordDocument") in head:
        return ".doc"
    if utf16("Workbook") in head or utf16("Book") in head:
        return ".xls"
    if utf16("PowerPoint Document") in head:
        return ".ppt"
    # Other things use OLE2 too — MicroStation .dgn drawings among them. Guessing ".doc"
    # sent a CAD file to LibreOffice and reported the refusal as a conversion failure,
    # which describes the wrong problem. An unknown container says so.
    return ""


def is_container_part(member):
    """Is this archive member the container's own plumbing rather than a document?

    A signed container — `.asice` here, `.edoc` elsewhere — is ASiC-E: a ZIP whose standard
    reserves `mimetype` at the root and everything under `META-INF/` for the format itself,
    the media type and the XAdES signature parts. They carry no tender content.

    On a deeply nested signed container most of the entries produced are this plumbing
    rather than documents. Left in, a signed tender pads the reading packet with signature
    XML and inflates every document count on the way. Taken out, they are still listed — as
    packaging, never dropped — because "nothing vanishes" is the claim this file keeps.

    This is not a judgement about usefulness, which stays forbidden. It is the format
    stating which of its own members are structure, the same way a ZIP's central directory
    is not a file.
    """
    return member == "mimetype" or member.startswith("META-INF/")


def packaging(rel_stem, path):
    """Record a container's own part: present, named, and not offered to a reader."""
    try:
        size = os.path.getsize(path)
    except OSError:
        size = -1
    return [{"source": rel_stem, "kind": "packaging", "bytes": size,
             "note": "container metadata, not a document", "preferred_for_agent": False}]


def unreadable(rel_stem, path, note):
    """Record a file no decoder could turn into text.

    Name, size and digest are kept because this list is the guarantee: anything absent from
    the extracted text is either here, named, or it is a bug. A count alone would let a
    missing specification hide behind a pile of CAD drawings.
    """
    try:
        size = os.path.getsize(path)
        h = hashlib.sha256()
        with open(path, "rb") as fh:
            for block in iter(lambda: fh.read(1 << 20), b""):
                h.update(block)
        digest = h.hexdigest()
    except OSError:
        size, digest = -1, ""
    # `on_disk` is where the bytes actually are, which is not derivable from `source` once a
    # file came out of an archive. Most files reported as scans live inside archives and
    # have no downloaded original to point at, so without this the fallback lane cannot
    # reach them at all. Naming the path is what makes them reachable.
    return [{"source": rel_stem, "kind": "unsupported", "note": note,
             "bytes": size, "sha256": digest, "on_disk": path,
             "preferred_for_agent": False}]


# --------------------------------------------------------------------------- converters
def pdf_to_md(path, media_dir, stem):
    """Every character in the text layer, page by page. Nothing detected, nothing inferred."""
    import fitz                                          # PyMuPDF
    doc = fitz.open(path)
    parts, media = [], []
    for number, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            parts.append("## Lapa %d\n\n%s" % (number, text))
        if not WITH_IMAGES:
            continue
        page_area = abs(page.rect.width * page.rect.height) or 1.0
        for img_i, info in enumerate(page.get_images(full=True), 1):
            xref = info[0]
            try:
                rects = page.get_image_rects(xref)
            except Exception:
                rects = []
            shown = max((abs(r.width * r.height) for r in rects), default=0.0)
            pix = fitz.Pixmap(doc, xref)
            if (pix.width < MIN_IMAGE_PX or pix.height < MIN_IMAGE_PX
                    or shown / page_area < MIN_IMAGE_PAGE_FRACTION):
                pix = None
                continue
            os.makedirs(media_dir, exist_ok=True)
            name = "%s_p%03d_%02d.png" % (stem, number, img_i)
            if pix.n - pix.alpha > 3:
                pix = fitz.Pixmap(fitz.csRGB, pix)
            pix.save(os.path.join(media_dir, name))
            pix = None
            media.append(name)
            parts.append("![%s](media/%s)" % (name, name))
    doc.close()
    return "\n\n".join(parts).strip(), media


_DOCX_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def _docx_xml_paragraphs(source, member_pattern, inside=None):
    """Paragraph texts from raw XML parts that python-docx never surfaces.

    What hides there is ordinary tender content: an object's street address in a header,
    the reference binding a qualification form to a clause of the nolikums in a footnote,
    page footers carrying the procurement ID. Parts are read in sorted order so the output
    is stable.
    """
    import xml.etree.ElementTree as ET
    lines = []
    try:
        with zipfile.ZipFile(source) as z:
            for member in sorted(n for n in z.namelist() if re.fullmatch(member_pattern, n)):
                try:
                    root = ET.fromstring(z.read(member))
                except ET.ParseError:
                    continue
                scopes = root.iter(_DOCX_W + inside) if inside else (root,)
                for scope in scopes:
                    for para in scope.iter(_DOCX_W + "p"):
                        text = "".join(t.text or "" for t in para.iter(_DOCX_W + "t")).strip()
                        if text:
                            lines.append(text)
    except (zipfile.BadZipFile, OSError):
        pass
    return lines


W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def _docx_numbering(path):
    """The numbering definitions Word keeps out of the text, read straight from the package.

    WHY THIS IS A DUMP AND NOT A RENDERER. Printing a Word number means implementing lvlText
    templates, formats, restarts and level overrides — a real tail of edge cases. Nobody
    downstream needs the printed number from us; they need enough to compute the one they
    want, and to know when they cannot. So this hands over `start`, `numFmt`, `lvlText`,
    `lvlRestart` and any `startOverride`, and prints nothing.
    """
    import xml.etree.ElementTree as ET
    try:
        with zipfile.ZipFile(path) as z:
            xml = z.read("word/numbering.xml")
    except (KeyError, zipfile.BadZipFile, OSError):
        return {}
    try:
        root = ET.fromstring(xml)
    except ET.ParseError:
        return {}

    def val(node, tag):
        el = node.find(W + tag)
        return el.get(W + "val") if el is not None else None

    abstracts = {}
    for a in root.findall(W + "abstractNum"):
        levels = {}
        for lvl in a.findall(W + "lvl"):
            levels[lvl.get(W + "ilvl")] = {
                "start": val(lvl, "start"),
                "numFmt": val(lvl, "numFmt"),
                "lvlText": val(lvl, "lvlText"),
                "lvlRestart": val(lvl, "lvlRestart"),
            }
        abstracts[a.get(W + "abstractNumId")] = levels

    numbering = {}
    for n in root.findall(W + "num"):
        abstract = val(n, "abstractNumId")
        overrides = {}
        for ov in n.findall(W + "lvlOverride"):
            start = ov.find(W + "startOverride")
            if start is not None:
                overrides[ov.get(W + "ilvl")] = start.get(W + "val")
        numbering[n.get(W + "numId")] = {
            "abstract": abstract,
            "levels": abstracts.get(abstract, {}),
            "overrides": overrides,
        }
    return numbering


def _numpr_of(pPr):
    """`(numId, ilvl)` out of one `w:pPr`, or `(None, None)`. Absent means absent."""
    if pPr is None:
        return None, None
    numPr = pPr.find(W + "numPr")
    if numPr is None:
        return None, None

    def val(tag):
        el = numPr.find(W + tag)
        return el.get(W + "val") if el is not None else None

    # ilvl stays None when Word did not write one. It used to be filled with 0, which is a
    # different claim: "this paragraph is at the top level" rather than "the file does not say".
    # A consumer computing a clause number from a fabricated depth gets a confident wrong
    # answer, which is the one outcome this whole sidecar exists to avoid.
    return val("numId"), val("ilvl")


def _style_numbering(style, styles):
    """The numbering a paragraph inherits from its style, following `w:basedOn` upwards.

    WHY THIS BRANCH EXISTS. Word resolves a paragraph's number in order: the paragraph's own
    `numPr`; failing that its style's; failing that the style it is based on, and so on up.
    A template that numbers through styles leaves every paragraph bare, and reading only the
    paragraph reports a document with no numbering at all — while the text cites its own
    clauses by number and the styles carry the depth spelled into their names.

    Returns `(numId, ilvl, source)` where source names the style the numbering came from, so a
    consumer can tell an inherited number from a paragraph's own.
    """
    element, seen = getattr(style, "_element", None), set()
    while element is not None and len(seen) < 12:
        style_id = element.get(W + "styleId")
        if style_id in seen:
            break                                   # a template can point a style at itself
        seen.add(style_id)
        num_id, ilvl = _numpr_of(element.find(W + "pPr"))
        if num_id is not None:
            name = element.find(W + "name")
            return num_id, ilvl, (name.get(W + "val") if name is not None else style_id)
        based = element.find(W + "basedOn")
        element = _style_by_id(styles, based.get(W + "val")) if based is not None else None
    return None, None, None


def _style_by_id(styles, style_id):
    """One `w:style` element out of the document's styles part, by id.

    Read off the XML rather than through python-docx's style collection: `basedOn` names an id,
    the collection is keyed by display name, and the two are not the same string.
    """
    if styles is None or not style_id:
        return None
    for element in styles.findall(W + "style"):
        if element.get(W + "styleId") == style_id:
            return element
    return None


def _paragraph_numbering(p, styles=None):
    """`(numId, ilvl, source)` for one paragraph, or `(None, None, None)`.

    This is the property flat markdown destroys: Word keeps a clause's number as a property of
    the paragraph or of its style, not as text, so `5.6.` never appears in the characters the
    converter reads.
    """
    num_id, ilvl = _numpr_of(p._p.find(W + "pPr"))
    if num_id is not None:
        return num_id, ilvl, "paragraph"
    # An explicit numId of 0 means Word switched numbering OFF for this paragraph, and it
    # overrides whatever the style says — so it returns above, as the paragraph's own answer,
    # and only a genuinely absent numPr falls through to the style.
    num_id, ilvl, from_style = _style_numbering(getattr(p, "style", None), styles)
    if num_id is not None:
        return num_id, ilvl, "style:%s" % (from_style or "?")
    return None, None, None


def docx_to_md(path, media_dir, stem):
    import docx
    d = docx.Document(path)
    # The styles part, so a paragraph that carries no numbering of its own can be resolved
    # through the style it uses and the styles that one is based on.
    try:
        styles_el = d.styles.element
    except Exception:
        styles_el = None
    parts = []
    numbered = []
    for p in d.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        level = 0
        if p.style is not None and (p.style.name or "").startswith("Heading"):
            digits = re.findall(r"\d+", p.style.name)
            level = int(digits[0]) if digits else 1
        # The index is the ordinal among the paragraphs that actually reach the markdown, not
        # among all of them: the empty ones are skipped just above, and a consumer counting
        # lines in the file would otherwise drift by exactly that many. The digest is there so
        # it can prove the alignment rather than assume it.
        num_id, ilvl, from_where = _paragraph_numbering(p, styles_el)
        if num_id is not None:
            numbered.append({
                "index": len(parts),
                "numId": num_id,
                "ilvl": int(ilvl) if (ilvl or "").isdigit() else None,
                "style": (p.style.name if p.style is not None else None),
                "from": from_where,
                "digest": hashlib.sha1(text.encode("utf-8")).hexdigest()[:8],
            })
        parts.append(("#" * min(level, 6) + " " + text) if level else text)
    for t_i, table in enumerate(d.tables, 1):
        rows = [[c.text.replace("\n", " ").strip() for c in row.cells] for row in table.rows]
        rows = [r for r in rows if any(r)]
        if not rows:
            continue
        if len(rows) == 1:
            # These forms carry approval stamps, column captions and whole declarations as
            # single-row tables; `< 2 rows` dropped their text entirely.
            parts.append("### Tabula %d\n\n| %s |" % (t_i, " | ".join(rows[0])))
            continue
        head, body = rows[0], rows[1:]
        parts.append("### Tabula %d\n\n| %s |\n| %s |\n%s"
                     % (t_i, " | ".join(head), " --- |" * len(head),
                        "\n".join("| %s |" % " | ".join(r) for r in body)))

    # Body text living inside drawing-canvas text boxes, then headers/footers, then foot-
    # and endnotes. Word repeats one header across sections, so lines are deduplicated.
    for title, pattern, inside in (
            (None, r"word/document\.xml", "txbxContent"),
            ("Galvenes un kājenes", r"word/(?:header|footer)\d*\.xml", None),
            ("Vēres", r"word/(?:footnotes|endnotes)\.xml", None)):
        seen, keep = set(), []
        for line in _docx_xml_paragraphs(path, pattern, inside):
            if line not in seen:
                seen.add(line)
                keep.append(line)
        if keep:
            parts.append(("### %s\n\n" % title if title else "") + "\n".join(keep))
    media = []
    for rel in d.part.rels.values():
        if "image" not in rel.reltype:
            continue
        blob = rel.target_part.blob
        if len(blob) < 4096:
            continue
        os.makedirs(media_dir, exist_ok=True)
        name = "%s_%s" % (stem, slug(os.path.basename(rel.target_part.partname)))
        with open(os.path.join(media_dir, name), "wb") as fh:
            fh.write(blob)
        media.append(name)
    if media:
        parts.append("\n".join("![%s](media/%s)" % (m, m) for m in media))
    structure = None
    if numbered:
        structure = {
            "schema": "structure/1",
            "note": ("paragraph facts as Word stores them; no number is printed here, and a "
                     "number computed from these is derived and must be labelled so. "
                     "`numId` is always the file's own value — `\"0\"` means Word switched "
                     "numbering off for that paragraph and overrides its style. `ilvl` is null "
                     "when the file does not say, never a filled-in zero. `from` is `paragraph` "
                     "or `style:<name>` — Word resolves a number from the paragraph first, then "
                     "its style, then the styles that one is based on."),
            "paragraphs": numbered,
            "numbering": {k: v for k, v in _docx_numbering(path).items()
                          if k in {n["numId"] for n in numbered}},
        }
    return "\n\n".join(parts).strip(), media, structure


def xlsx_to_md(path):
    import openpyxl
    # A formula cell holds a computed value only if the producing program cached one.
    # Excel always caches; other generators may not, and data_only=True then yields None —
    # a silently empty cell in the middle of a bill of quantities. Depending on every
    # producer having cached one is luck rather than a guarantee, so a second pass with
    # formulas visible supplies the formula text as the deterministic fallback.
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    wb_f = openpyxl.load_workbook(path, read_only=True, data_only=False)
    parts = []
    for ws, ws_f in zip(wb.worksheets, wb_f.worksheets):
        rows = []
        for row, row_f in zip(ws.iter_rows(values_only=True), ws_f.iter_rows(values_only=True)):
            cells = []
            for c, f in zip(row, row_f):
                if c is None and isinstance(f, str) and f.startswith("="):
                    c = f
                cells.append("" if c is None
                             else str(c).replace("\n", " ").replace("|", "\\|").strip())
            while cells and not cells[-1]:
                cells.pop()
            if any(cells):
                rows.append(cells)
        if not rows:
            continue
        width = max(len(r) for r in rows)
        rows = [r + [""] * (width - len(r)) for r in rows]
        head, body = rows[0], rows[1:]
        parts.append("## %s\n\n| %s |\n| %s |\n%s"
                     % (ws.title, " | ".join(head), " --- |" * width,
                        "\n".join("| %s |" % " | ".join(r) for r in body)))
    wb.close()
    wb_f.close()
    return "\n\n".join(parts).strip(), []


def pptx_to_md(path, media_dir, stem):
    from pptx import Presentation
    prs = Presentation(path)
    parts, media = [], []
    for number, slide in enumerate(prs.slides, 1):
        lines = ["## Slaids %d" % number]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                rows = [[c.text.replace("\n", " ").strip() for c in row.cells]
                        for row in shape.table.rows]
                rows = [r for r in rows if any(r)]
                if len(rows) == 1:
                    lines.append("| %s |" % " | ".join(rows[0]))
                elif rows:
                    lines.append("| %s |\n| %s |\n%s"
                                 % (" | ".join(rows[0]), " --- |" * len(rows[0]),
                                    "\n".join("| %s |" % " | ".join(r) for r in rows[1:])))
            if shape.shape_type == 13 and getattr(shape, "image", None):     # PICTURE
                blob = shape.image.blob
                if len(blob) < 4096:
                    continue
                os.makedirs(media_dir, exist_ok=True)
                name = "%s_s%03d.%s" % (stem, number, shape.image.ext)
                with open(os.path.join(media_dir, name), "wb") as fh:
                    fh.write(blob)
                media.append(name)
                lines.append("![%s](media/%s)" % (name, name))
        if len(lines) > 1:
            parts.append("\n\n".join(lines))
    return "\n\n".join(parts).strip(), media


def ifc_to_md(path):
    """A BIM model is ISO-10303-21 text, not a binary, and calling it unreadable was wrong.

    The geometry is numeric and carries no words; the words live in the quoted string
    attributes — element names, property-set names, classifications. On a real 283 MB model
    that is a great many distinct strings recovered in about a second, among them fire
    resistance classes and their ratings. Strings are deduplicated and sorted so the output
    is stable.
    """
    quoted = re.compile(r"'([^']*)'")
    x2 = re.compile(r"\\X2\\((?:[0-9A-Fa-f]{4})+)\\X0\\")
    x1 = re.compile(r"\\X\\([0-9A-Fa-f]{2})")
    guid = re.compile(r"^[0-9A-Za-z_$]{22}$")           # an IfcGloballyUniqueId, never prose

    def decode(s):
        s = x2.sub(lambda m: bytes.fromhex(m.group(1)).decode("utf-16-be", "replace"), s)
        return x1.sub(lambda m: chr(int(m.group(1), 16)), s)

    seen = set()
    with open(path, encoding="utf-8", errors="replace") as fh:
        for line in fh:
            for raw in quoted.findall(line):
                if not raw or guid.match(raw):
                    continue
                text = decode(raw).strip()
                if text and any(c.isalpha() for c in text):
                    seen.add(text)
    if not seen:
        return "", []
    return ("## Modeļa apzīmējumi (%d unikāli)\n\n" % len(seen)
            + "\n".join("- %s" % s for s in sorted(seen))), []


def textual_to_md(path, kind):
    with open(path, "rb") as fh:
        blob = fh.read(MAX_MEMBER_BYTES)
    # Strict UTF-8 first; on failure, windows-1257, the Baltic legacy codepage. A bare
    # "replace" here turned Būvdarbu into B�vdarbu — diacritics destroyed with no record
    # of the damage, in files whose whole value is their words.
    try:
        raw = blob.decode("utf-8-sig")
    except UnicodeDecodeError:
        raw = blob.decode("windows-1257", "replace")
    if kind in ("csv", "tsv"):
        import csv
        delim = "\t" if kind == "tsv" else ","
        rows = [r for r in csv.reader(io.StringIO(raw), delimiter=delim) if any(r)]
        if len(rows) >= 2:
            width = max(len(r) for r in rows)
            rows = [[c.replace("|", "\\|").strip() for c in r] + [""] * (width - len(r))
                    for r in rows]
            return ("| %s |\n| %s |\n%s"
                    % (" | ".join(rows[0]), " --- |" * width,
                       "\n".join("| %s |" % " | ".join(r) for r in rows[1:])), [])
    if kind == "markdown":
        return raw.strip(), []
    return "```%s\n%s\n```" % ("" if kind == "text" else kind, raw.strip()), []


# ------------------------------------------------------------------------ legacy formats
# Buyers still attach files saved by Word 97. Those carry ordinary tender text — one of them
# held the procurement commission's answers to bidders — and leaving them unread is exactly
# the silent gap this module exists to prevent. LibreOffice converts them without a model
# and without a network call, so the guarantee of a deterministic, token-free extract holds.
#
# The OpenDocument family joins them for the same reason and by the same road: the converter
# is already installed, already deterministic and already offline, so reading an .ods is one
# entry in this table rather than a second dependency.
LEGACY = {".doc": "docx", ".xls": "xlsx", ".ppt": "pptx", ".rtf": "docx",
          ".ods": "xlsx", ".odt": "docx", ".odp": "pptx"}


def _soffice():
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    return None


def legacy_to_md(path, media_dir, stem, kind):
    import subprocess
    import tempfile
    binary = _soffice()
    if not binary:
        raise RuntimeError("LibreOffice is not installed, so %s cannot be read" % kind)
    with tempfile.TemporaryDirectory() as tmp:
        # A private profile per call: LibreOffice refuses to run twice against one profile,
        # and a shared profile would make the result depend on run order.
        proc = subprocess.run(
            [binary, "--headless", "--norestore",
             "-env:UserInstallation=file://%s" % os.path.join(tmp, "profile").replace("\\", "/"),
             "--convert-to", LEGACY[kind], "--outdir", tmp,
             plainpath(os.path.abspath(path))],
            stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=180)
        produced = [os.path.join(tmp, n) for n in os.listdir(tmp)
                    if n.lower().endswith("." + LEGACY[kind])]
        if not produced:
            raise RuntimeError("LibreOffice produced nothing: %s"
                               % proc.stdout.decode("utf-8", "replace").strip()[:120])
        modern = produced[0]
        if LEGACY[kind] == "docx":
            return docx_to_md(modern, media_dir, stem)
        if LEGACY[kind] == "xlsx":
            return xlsx_to_md(modern)
        return pptx_to_md(modern, media_dir, stem)


# ----------------------------------------------------------------------------- unpacking
def unpack(path, dest, kind, budget):
    """Extract an archive, refusing traversal and honouring the expansion budget.

    The destination is always emptied first. py7zr raises FileExistsError when it extracts
    into a directory that already holds those names, and that error surfaces as
    `archive not unpacked` — leaving the archive in the output looking present and
    unreadable, which on a building project is most of the tender.
    """
    if os.path.isdir(dest):
        shutil.rmtree(dest, ignore_errors=True)
    os.makedirs(dest, exist_ok=True)
    if kind == ".zip":
        with zipfile.ZipFile(path) as z:
            members = [m for m in z.infolist() if not m.filename.endswith("/")]
            if len(members) > MAX_MEMBERS:
                raise ValueError("archive holds %d members, limit %d" % (len(members), MAX_MEMBERS))
            for m in members:
                name = m.filename.replace("\\", "/")
                if name.startswith("/") or ".." in name.split("/") or os.path.isabs(name):
                    raise ValueError("unsafe path in archive: %r" % name)
                budget[0] -= m.file_size
                if budget[0] < 0:
                    raise ValueError("archive expands past the size limit")
                target = os.path.join(dest, *[slug(p) for p in name.split("/") if p])
                os.makedirs(os.path.dirname(target), exist_ok=True)
                with z.open(m) as src, open(target, "wb") as dst:
                    shutil.copyfileobj(src, dst)
        return
    binary = shutil.which("7z") or shutil.which("7za")
    if binary:
        return _unpack_7z_cli(binary, path, dest, budget)
    # Fallback where the CLI is absent (a local Windows checkout). On the runner the CLI
    # is installed and 20 seconds faster on the 484 MB acceptance archive, and using it
    # removes one library-version axis from the "same bytes, same text" claim.
    import py7zr
    with py7zr.SevenZipFile(path, "r") as z:
        names = z.getnames()
        if len(names) > MAX_MEMBERS:
            raise ValueError("archive holds %d members, limit %d" % (len(names), MAX_MEMBERS))
        for name in names:
            clean = name.replace("\\", "/")
            if clean.startswith("/") or ".." in clean.split("/") or os.path.isabs(clean):
                raise ValueError("unsafe path in archive: %r" % name)
        total = sum(getattr(i, "uncompressed", 0) or 0 for i in z.list())
        budget[0] -= total
        if budget[0] < 0:
            raise ValueError("archive expands to %s bytes, past the limit" % f"{total:,}")
        z.extractall(dest)


def _unpack_7z_cli(binary, path, dest, budget):
    """Same traversal and budget rules as the zip path, enforced from `7z l -slt` output."""
    import subprocess
    proc = subprocess.run([binary, "l", "-slt", "--", os.path.abspath(path)],
                          stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=600)
    listing = proc.stdout.decode("utf-8", "replace")
    if proc.returncode != 0:
        raise ValueError("7z could not list the archive: %s" % listing[-160:])
    # Entries follow the "----------" separator as "Field = value" blocks; the block before
    # it describes the archive itself.
    names, total, current = [], 0, {}
    body = listing.split("\n----------\n", 1)
    blocks = body[1] if len(body) == 2 else ""
    for line in blocks.splitlines() + [""]:
        if line.strip():
            key, _, value = line.partition(" = ")
            current[key.strip()] = value.strip()
            continue
        if current:
            if "D" not in current.get("Attributes", ""):
                names.append(current.get("Path", ""))
                try:
                    total += int(current.get("Size") or 0)
                except ValueError:
                    pass
            current = {}
    if len(names) > MAX_MEMBERS:
        raise ValueError("archive holds %d members, limit %d" % (len(names), MAX_MEMBERS))
    for name in names:
        clean = name.replace("\\", "/")
        if clean.startswith("/") or ".." in clean.split("/") or os.path.isabs(clean):
            raise ValueError("unsafe path in archive: %r" % name)
    budget[0] -= total
    if budget[0] < 0:
        raise ValueError("archive expands to %s bytes, past the limit" % f"{total:,}")
    proc = subprocess.run([binary, "x", "-y", "-o" + plainpath(os.path.abspath(dest)), "--",
                           plainpath(os.path.abspath(path))],
                          stdout=subprocess.PIPE, stderr=subprocess.STDOUT, timeout=1800)
    if proc.returncode != 0:
        raise ValueError("7z extraction failed: %s"
                         % proc.stdout.decode("utf-8", "replace")[-160:])


# -------------------------------------------------------------------------------- driver
def convert(path, out_root, rel_stem, budget, depth=0):
    """One file to Markdown, recursing into archives. Returns manifest entries."""
    kind = sniff(path)
    stem = slug(os.path.splitext(os.path.basename(path))[0])
    # `rel_stem` is a logical address joined with forward slashes, so this is where it
    # becomes a real path and where both of the prefix's rules have to be applied.
    doc_dir = fspath(os.path.join(out_root, rel_stem))
    media_dir = os.path.join(doc_dir, "media")

    if kind in ARCHIVES:
        if depth >= MAX_DEPTH:
            return unreadable(rel_stem, path, "archive nested deeper than %d" % MAX_DEPTH)
        # Scratch space, addressed by identity rather than by name. Naming it after the
        # archive put the destination directory at the exact path of the archive file for
        # any NESTED archive — rel_stem keeps the extension below the top level — so
        # makedirs() hit the file and raised "File exists". That is how a nested archive
        # holding the whole building project stayed unread while the run reported success.
        # A hash also keeps the path short, which matters when unpacking deep trees.
        inner = fspath(os.path.join(out_root, "_unpacked",
                                    hashlib.sha1(rel_stem.encode("utf-8")).hexdigest()[:16]))
        try:
            unpack(path, inner, kind, budget)
        except Exception as exc:
            return unreadable(rel_stem, path, "archive not unpacked: %s" % str(exc)[:120])
        out = []
        # dirs.sort() pins the walk order: os.walk otherwise follows filesystem order,
        # which would make entry order — and any diff of two runs — depend on the disk.
        for root, dirs, names in os.walk(inner):
            dirs.sort()
            for name in sorted(names):
                child = os.path.join(root, name)
                member = os.path.relpath(child, inner).replace("\\", "/")
                child_rel = os.path.join(rel_stem, member).replace("\\", "/")
                if is_container_part(member):
                    out.extend(packaging(child_rel, child))
                    continue
                out.extend(convert(child, out_root, child_rel, budget, depth + 1))
        return out

    # Word is the only format that hides structure outside the text, so it is the only one that
    # returns a third value. `(result + (None,))[:3]` keeps every other converter's signature
    # exactly as it was rather than making them all carry a field they have nothing to put in.
    structure = None
    try:
        if kind == ".pdf":
            md, media = pdf_to_md(path, media_dir, stem)
        elif kind == ".docx":
            md, media, structure = docx_to_md(path, media_dir, stem)
        elif kind == ".xlsx":
            md, media = xlsx_to_md(path)
        elif kind == ".pptx":
            md, media = pptx_to_md(path, media_dir, stem)
        elif kind == ".ifc":
            md, media = ifc_to_md(path)
        elif kind in LEGACY:
            md, media, structure = (legacy_to_md(path, media_dir, stem, kind) + (None,))[:3]
        elif kind in TEXTUAL:
            md, media = textual_to_md(path, TEXTUAL[kind])
        else:
            return unreadable(rel_stem, path,
                              "no deterministic text extraction for %s" % (kind or "unknown"))
    except Exception as exc:
        return unreadable(rel_stem, path, "conversion failed: %s" % str(exc)[:140])

    if len(md.strip()) < MIN_USEFUL_CHARS:
        return unreadable(rel_stem, path, "no text layer — a scan or a vector drawing")

    os.makedirs(doc_dir, exist_ok=True)
    md_path = os.path.join(doc_dir, "document.md")
    with open(md_path, "w", encoding="utf-8") as fh:
        fh.write("# %s\n\n%s\n" % (os.path.basename(path), md))
    entry = {"source": rel_stem,
             "kind": "markdown+media" if media else "markdown",
             "markdown_path": os.path.relpath(md_path, out_root).replace("\\", "/"),
             "markdown_chars": len(md),
             "media_count": len(media),
             "preferred_for_agent": True}
    # Beside the text, never inside it: document.md is byte for byte what it was before this
    # existed, so digests and any consumer's quote rule are untouched. A consumer that does
    # not know about the sidecar cannot notice it.
    if structure:
        structure_path = os.path.join(doc_dir, "structure.json")
        with open(structure_path, "w", encoding="utf-8") as fh:
            json.dump(dict(structure, source=rel_stem), fh, ensure_ascii=False)
        entry["structure_path"] = os.path.relpath(structure_path, out_root).replace("\\", "/")
    return [entry]


def read_markdown(path, offset_chars=0, max_chars=50000):
    """Chunked reading, so a 300,000-character specification arrives in usable pieces."""
    with open(path, encoding="utf-8") as fh:
        text = fh.read()
    chunk = text[offset_chars:offset_chars + max_chars]
    nxt = offset_chars + len(chunk)
    return {"markdown": chunk, "total_chars": len(text), "offset_chars": offset_chars,
            "has_more": nxt < len(text), "next_offset_chars": nxt if nxt < len(text) else None}


def main(argv=None):
    ap = argparse.ArgumentParser(description="Normalize downloaded tender documents to Markdown.")
    ap.add_argument("--in", dest="src", required=True, help="fetch output directory")
    ap.add_argument("--out", required=True, help="normalized output directory")
    ap.add_argument("--with-images", action="store_true",
                    help="also extract pictures from PDFs (off by default: they add no text)")
    ap.add_argument("--keep-unpacked", action="store_true",
                    help="keep the archive scratch so a later stage can read files that only "
                         "exist inside archives (the scan lane needs this)")
    args = ap.parse_args(argv)

    global WITH_IMAGES, KEEP_UNPACKED
    WITH_IMAGES = args.with_images
    KEEP_UNPACKED = args.keep_unpacked
    started = time.time()
    # Prefixed once, here, so that everything derived from them inherits it and
    # `os.path.relpath` still answers correctly between two paths of the same shape.
    src, out = fspath(args.src), fspath(args.out)
    with open(os.path.join(src, "manifest.json"), encoding="utf-8") as fh:
        manifest = json.load(fh)
    os.makedirs(out, exist_ok=True)
    budget = [MAX_EXPANDED]
    entries, converted = [], {}

    for record in manifest["documents"]:
        for f in record["files"]:
            # `f["path"]` is the manifest's logical address and uses forward slashes, so
            # this is the other place a logical address becomes a real one. Without the
            # normalisation the file is simply not found, `sniff` answers "unknown", and a
            # perfectly readable tender is reported as having no text extraction.
            path = fspath(os.path.join(src, f["path"]))
            # The same file can be listed by two records (current and archive sections).
            # Converting it twice wastes minutes and, for archives, fails the second time.
            if path in converted:
                for prior in converted[path]:
                    entries.append(dict(prior, record_id=record["id"],
                                        record_title=record["title"],
                                        section=record["section"],
                                        also_listed_under=prior.get("record_id")))
                continue
            rel = os.path.join(record["section"], slug(os.path.splitext(f["filename"])[0]))
            produced = []
            for entry in convert(path, out, rel.replace("\\", "/"), budget):
                entry.update({"record_id": record["id"], "record_title": record["title"],
                              "section": record["section"], "original_file": f["filename"],
                              "original_sha256": f["sha256"]})
                produced.append(entry)
                entries.append(entry)
            converted[path] = produced

            # One line said only "unsupported", and the reason lived in a manifest inside a
            # 465 MB artifact — so diagnosing a silently unread document meant downloading
            # half a gigabyte. Reasons are printed here, deduplicated so 156 CAD drawings
            # cost one line rather than 156.
            tally = {}
            for e in produced:
                tally[e["kind"]] = tally.get(e["kind"], 0) + 1
            print("  %-58s %s" % (f["filename"][:58],
                                  ", ".join("%d %s" % (n, k) for k, n in sorted(tally.items()))),
                  flush=True)
            notes = {}
            for e in produced:
                if e["kind"] == "unsupported":
                    notes.setdefault(e.get("note", "no reason recorded"), []).append(e["source"])
            for note, sources in notes.items():
                print("      ! %s  (%d file(s), e.g. %s)"
                      % (note, len(sources), sources[0][-70:]), flush=True)

    # The audit list, deduplicated by digest so the same drawing reached through two records
    # is one gap and not two. This is the claim "nothing useful was lost" made checkable:
    # every file not represented in the text is named here.
    gaps, seen_digest = [], set()
    for e in entries:
        if e["kind"] != "unsupported":
            continue
        key = e.get("sha256") or e["source"]
        if key in seen_digest:
            continue
        seen_digest.add(key)
        gap = {"file": e["source"], "bytes": e.get("bytes", -1),
               "sha256": e.get("sha256", ""), "reason": e.get("note", "")}
        # Where the bytes are, relative to the pack, so a later stage can open the file
        # without re-deriving it. Only meaningful while the unpacked tree survives.
        if KEEP_UNPACKED and e.get("on_disk"):
            try:
                gap["path"] = os.path.relpath(e["on_disk"], src).replace("\\", "/")
            except ValueError:
                pass
        gaps.append(gap)
    gaps.sort(key=lambda g: -g["bytes"])

    # Cleared last, and only when nobody asked to keep it. The scratch holds the ONLY copy of
    # every file that came out of an archive, so removing it before the fallback lane has run
    # is what limited that lane to 6 of 26 scans on a measured day.
    if not KEEP_UNPACKED:
        shutil.rmtree(os.path.join(out, "_unpacked"), ignore_errors=True)

    elapsed = time.time() - started
    # `entries` maps records to files and legitimately lists a file under every record
    # that carries it; the headline numbers must not. Counting duplicates inflates the
    # character total — precision theater in the one summary everyone reads.
    unique = [e for e in entries if not e.get("also_listed_under")]
    doc = {"schema": 2, "procurement_id": manifest["procurement_id"],
           "entries": len(entries),
           "markdown": sum(1 for e in unique if e.get("preferred_for_agent")),
           "unsupported": sum(1 for e in unique if e["kind"] == "unsupported"),
           "chars": sum(e.get("markdown_chars", 0) for e in unique),
           "seconds": round(elapsed, 1),
           "images_extracted": WITH_IMAGES,
           "unreadable_files": gaps,
           "documents": entries}
    with open(os.path.join(out, "manifest_normalized.json"), "w", encoding="utf-8") as fh:
        json.dump(doc, fh, ensure_ascii=False, indent=2)

    archive = os.path.join(src, "eis_%s_normalized.zip" % manifest["procurement_id"])
    with zipfile.ZipFile(archive, "w", zipfile.ZIP_DEFLATED, compresslevel=6,
                         allowZip64=True) as z:
        for root, dirs, names in os.walk(out):
            dirs.sort()
            for name in sorted(names):
                full = os.path.join(root, name)
                z.write(full, os.path.relpath(full, out).replace("\\", "/"))

    print("normalized %d entries · %d markdown · %s chars · %.1fs"
          % (doc["entries"], doc["markdown"], f"{doc['chars']:,}", elapsed))
    if gaps:
        print("%d file(s) yielded no text — every one of them named here:" % len(gaps))
        for g in gaps[:20]:
            print("   %10s B  %-58s %s"
                  % (f"{g['bytes']:,}", g["file"][-58:], g["reason"][:60]))
        if len(gaps) > 20:
            print("   ... %d more, all listed in manifest_normalized.json"
                  % (len(gaps) - 20))
    return 0


if __name__ == "__main__":
    import sys
    sys.exit(main())
